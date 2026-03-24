"""
Generates claude/presentation.pptx — De Slimste Mens Timer project presentation.
Run: python claude/make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy, os

# ── Palette ──────────────────────────────────────────────────────────────────
RED    = RGBColor(0xB7, 0x1C, 0x1C)
ORANGE = RGBColor(0xFF, 0xA7, 0x26)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
CARD   = RGBColor(0x8B, 0x14, 0x14)   # slightly lighter red for cards
DIM    = RGBColor(0xFF, 0xD0, 0x80)   # pale orange for body text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_background(slide, color=RED):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        9,  # oval
        left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def txb(slide, text, left, top, width, height,
        font_size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb

def heading(slide, text, top=Inches(0.45), size=28):
    # Orange heading bar
    add_rect(slide, 0, top - Inches(0.05), SLIDE_W, Inches(0.55),
             RGBColor(0x00, 0x00, 0x00))
    bar = add_rect(slide, 0, top - Inches(0.05), Inches(0.06), Inches(0.55), ORANGE)
    txb(slide, text,
        Inches(0.25), top, SLIDE_W - Inches(0.5), Inches(0.5),
        font_size=size, bold=True, color=ORANGE)

def bullet_list(slide, items, left, top, width, height,
                font_size=17, color=WHITE, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return tb

def code_block(slide, code, left, top, width, height, font_size=13):
    add_rect(slide, left, top, width, height, RGBColor(0x1A, 0x1A, 0x1A))
    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xFF, 0xD0, 0x80)
        run.font.name = "Cascadia Code"

def card(slide, left, top, width, height):
    add_rect(slide, left, top, width, height, CARD)

def table_slide(slide, headers, rows, left, top, width, col_widths=None,
                header_size=14, row_size=13):
    """Draw a simple table using rectangles and text boxes."""
    row_h = Inches(0.38)
    if col_widths is None:
        col_widths = [width / len(headers)] * len(headers)

    # Header row
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_widths[i], row_h, RGBColor(0x7B, 0x10, 0x10))
        txb(slide, h, x + Inches(0.08), top + Inches(0.05),
            col_widths[i] - Inches(0.1), row_h - Inches(0.05),
            font_size=header_size, bold=True, color=ORANGE)
        x += col_widths[i]

    # Data rows
    for ri, row in enumerate(rows):
        y = top + row_h * (ri + 1)
        bg = RGBColor(0x9C, 0x18, 0x18) if ri % 2 == 0 else CARD
        x = left
        for ci, cell in enumerate(row):
            add_rect(slide, x, y, col_widths[ci], row_h, bg)
            txb(slide, cell, x + Inches(0.08), y + Inches(0.05),
                col_widths[ci] - Inches(0.1), row_h - Inches(0.05),
                font_size=row_size, color=WHITE)
            x += col_widths[ci]

def accent_line(slide, top):
    add_rect(slide, Inches(0.4), top, Inches(12.5), Inches(0.025), ORANGE)

# ── Slides ────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]   # completely blank

# ── 1. Title slide ────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)

# Decorative circles
for x, y, s, a in [
    (Inches(11.5), Inches(-0.8), Inches(3),   CARD),
    (Inches(12.2), Inches(5.5),  Inches(2.5), CARD),
    (Inches(-0.5), Inches(5.8),  Inches(2),   CARD),
    (Inches(0.2),  Inches(-0.3), Inches(1.5), CARD),
]:
    add_circle(sl, x, y, s, a)

# Orange accent bar on left
add_rect(sl, 0, Inches(2.0), Inches(0.18), Inches(3.5), ORANGE)

# Title
txb(sl, "De Slimste Mens Timer",
    Inches(0.5), Inches(2.0), Inches(10), Inches(1.4),
    font_size=52, bold=True, color=WHITE)

# Subtitle
txb(sl, "From solo Android app to multiplayer scoreboard",
    Inches(0.5), Inches(3.5), Inches(10), Inches(0.7),
    font_size=24, color=ORANGE)

# Bottom line
txb(sl, "Built entirely with Claude Code  ·  March 2026",
    Inches(0.5), Inches(6.5), Inches(10), Inches(0.5),
    font_size=16, color=DIM, italic=True)

# ── 2. What is it? ────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "What is it?")

txb(sl, "A digital replacement for the physical timer devices in the Dutch board game\n"
        "De Slimste Mens (The Smartest Person).",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.8),
    font_size=18, color=WHITE)

# Phase 1 card
card(sl, Inches(0.4), Inches(2.1), Inches(5.9), Inches(3.8))
txb(sl, "Phase 1 — Android App",
    Inches(0.6), Inches(2.2), Inches(5.5), Inches(0.5),
    font_size=18, bold=True, color=ORANGE)
bullet_list(sl, [
    "60-second countdown timer",
    "Start / Stop toggle",
    "+20 and −20 score buttons",
    "Hamburger menu (reset, set score, exit)",
    "Portrait-only, dark red / orange palette",
    "Works fully standalone — no server needed",
], Inches(0.6), Inches(2.75), Inches(5.5), Inches(3.0), font_size=16)

# Phase 2 card
card(sl, Inches(6.9), Inches(2.1), Inches(5.9), Inches(3.8))
txb(sl, "Phase 2 — Multiplayer Extension",
    Inches(7.1), Inches(2.2), Inches(5.5), Inches(0.5),
    font_size=18, bold=True, color=ORANGE)
bullet_list(sl, [
    "ASP.NET Core 8 web server",
    "Browser scoreboard — live score tiles",
    "Players join via QR code scan",
    "SignalR for real-time browser updates",
    "App pushes scores via plain HTTP",
    "Server connectivity is optional",
], Inches(7.1), Inches(2.75), Inches(5.5), Inches(3.0), font_size=16)

# ── 3. Original Prompt ────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "The Original Prompt")

add_rect(sl, Inches(0.4), Inches(1.15), Inches(0.06), Inches(5.5), ORANGE)

txb(sl,
    '"I want to build an android app built with Delphi / Firemonkey to be used '
    'in conjunction with the \'De Slimste Mens\' board game. The app will be a '
    'replacement for the board games own timer devices.\n\n'
    'The timer / scoring app is a simple app with a countdown timer that is '
    'initially set to 60 seconds. When pressing the start/stop button, the score '
    'starts counting down at 1 point per second. When pressed again, the timer '
    'stops counting down. There are 2 more buttons, marked +20 and -20, which '
    'add or subtract 20 seconds / points from the current score.\n\n'
    'The app will only have to have a portrait oriented layout. The layout of the '
    'screen is so that the score is displayed more prominently at the top, centered '
    'in the screen. The 3 buttons are below that. From left to right they are +20, '
    'start/stop, -20. The background is a slightly darker red. The score and text '
    'on the buttons is white. The buttons themselves are round, and slightly '
    'orange/yellow in color."',
    Inches(0.65), Inches(1.2), Inches(12.0), Inches(5.5),
    font_size=16, color=WHITE, italic=True)

# ── 4. Phase 1 — Steps 1–4 ───────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Phase 1 — Android App  ·  Steps 1–4")

steps = [
    ("1", "Project Scaffolding",
     "SlimsteMensTimer.dpr + .dproj · Android target · Portrait lock\nDebug + Release configs · empty MainFrm ready for UI"),
    ("2", "UI Layout",
     "rectBackground dark red #B71C1C · lblScore large white centered\nbtnMinus20, btnStartStop, btnPlus20 — round orange #FFA726"),
    ("3", "Timer Logic",
     "tmrCountdown: TTimer (1 s) · decrements score · clamps to 0\ntxtStartStopClick toggles enabled state + button label"),
    ("4", "Style Refinement",
     "Score font 120 pt → 150 pt · button sizing + spacing\nAnchors set for Android screen size responsiveness"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.2) + i * Inches(1.45)
    card(sl, Inches(0.4), y, Inches(12.4), Inches(1.3))
    add_rect(sl, Inches(0.4), y, Inches(0.7), Inches(1.3), ORANGE)
    txb(sl, num, Inches(0.4), y + Inches(0.2), Inches(0.7), Inches(0.8),
        font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, title, Inches(1.25), y + Inches(0.08), Inches(11.3), Inches(0.45),
        font_size=18, bold=True, color=ORANGE)
    txb(sl, desc, Inches(1.25), y + Inches(0.52), Inches(11.3), Inches(0.7),
        font_size=14, color=WHITE)

# ── 5. Phase 1 — Steps 5–7 ───────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Phase 1 — Android App  ·  Steps 5–7")

# Step 5
card(sl, Inches(0.4), Inches(1.2), Inches(12.4), Inches(1.6))
add_rect(sl, Inches(0.4), Inches(1.2), Inches(0.7), Inches(1.6), ORANGE)
txb(sl, "5", Inches(0.4), Inches(1.45), Inches(0.7), Inches(0.8),
    font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Hamburger Menu", Inches(1.25), Inches(1.28), Inches(11.3), Inches(0.45),
    font_size=18, bold=True, color=ORANGE)
txb(sl, "btnMenu (☰) top-left toggles pnlMenu panel  ·  Reset score (→ 60, stops timer)\n"
        "Score instellen: TDialogService.InputQuery, validates 0–1000  ·  Afsluiten: Application.Terminate",
    Inches(1.25), Inches(1.72), Inches(11.3), Inches(0.9), font_size=14, color=WHITE)

# Step 6
card(sl, Inches(0.4), Inches(3.0), Inches(12.4), Inches(1.9))
add_rect(sl, Inches(0.4), Inches(3.0), Inches(0.7), Inches(1.9), ORANGE)
txb(sl, "6", Inches(0.4), Inches(3.3), Inches(0.7), Inches(0.8),
    font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "ScoreManager Class", Inches(1.25), Inches(3.08), Inches(11.3), Inches(0.45),
    font_size=18, bold=True, color=ORANGE)
txb(sl, "IScoreManager interface · TScoreManager (TInterfacedObject)\n"
        "Increase (uncapped) · Decrease (clamps to 0) · SetScore (raises EArgumentOutOfRangeException)\n"
        "Reset → 60  ·  Decorator-ready: FScoreManager: IScoreManager swappable at runtime",
    Inches(1.25), Inches(3.52), Inches(11.3), Inches(1.3), font_size=14, color=WHITE)

# Step 7
card(sl, Inches(0.4), Inches(5.1), Inches(12.4), Inches(1.7))
add_rect(sl, Inches(0.4), Inches(5.1), Inches(0.7), Inches(1.7), ORANGE)
txb(sl, "7", Inches(0.4), Inches(5.35), Inches(0.7), Inches(0.8),
    font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "DUnitX Unit Tests — 20 tests", Inches(1.25), Inches(5.18), Inches(11.3), Inches(0.45),
    font_size=18, bold=True, color=ORANGE)
txb(sl, "Create · Reset · Increase · Decrease · SetScore (valid + boundary) · SetScore exceptions\n"
        "EArgumentOutOfRangeException asserted for <0 and >1000  ·  All 20 passing on Win32/Win64",
    Inches(1.25), Inches(5.62), Inches(11.3), Inches(1.0), font_size=14, color=WHITE)

# ── 6. Phase 2 — Multiplayer Prompt ──────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Phase 2 — The Multiplayer Prompt")

add_rect(sl, Inches(0.4), Inches(1.15), Inches(0.06), Inches(4.5), ORANGE)
txb(sl,
    '"Now we want to add a shared scoreboard. A web page that shows all connected '
    "players' scores, updated in real time. A server manages game sessions. Players "
    'join a session by scanning a QR code from the app.\n\n'
    'The existing app should continue to work offline. Server connectivity is optional. '
    'When connected, score changes should appear on the scoreboard within about one second.\n\n'
    'The server should be an ASP.NET Core application. Use SignalR for real-time browser '
    'updates. The Delphi app pushes scores via plain HTTP — no SignalR on the app side."',
    Inches(0.65), Inches(1.2), Inches(12.0), Inches(4.5),
    font_size=17, color=WHITE, italic=True)

txb(sl, "Reconstructed from project documentation — no session transcript was preserved.",
    Inches(0.5), Inches(6.5), Inches(12.0), Inches(0.5),
    font_size=13, color=DIM, italic=True)

# ── 7. Phase 2 — S1–S4 Server Foundation ─────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Phase 2 — Server  ·  S1–S4")

cols = [
    ("S1\nScaffolding",
     ["ASP.NET Core 8 Web API", "QRCoder + Swashbuckle NuGet", "SignalR + CORS + Swagger",
      "appsettings: port 5000, StaleSeconds=30", "Startup log prints LAN IP"]),
    ("S2\nModels + Store",
     ["SessionState: Lobby → Active → Ended", "Session, Player models",
      "SessionStore singleton (ConcurrentDictionary)", "6-char codes, unambiguous charset",
      "Thread-safe CRUD operations"]),
    ("S3\nREST API",
     ["POST /api/sessions — create", "POST /players — register", "PUT /score — push score",
      "POST /heartbeat — keep-alive", "GET /qr + /appqr — QR PNGs"]),
    ("S4\nSignalR Hub",
     ["GameHub.JoinSessionGroup(id)", "PlayerJoined / ScoreUpdated",
      "PlayerWentStale / PlayerReturned", "GameStarted / GameEnded",
      "IHubContext wired to controllers"]),
]
col_w = Inches(3.1)
for i, (title, bullets) in enumerate(cols):
    x = Inches(0.3) + i * (col_w + Inches(0.1))
    card(sl, x, Inches(1.15), col_w, Inches(5.9))
    txb(sl, title, x + Inches(0.1), Inches(1.2), col_w - Inches(0.2), Inches(0.7),
        font_size=16, bold=True, color=ORANGE)
    accent_line(sl, Inches(1.92))
    bullet_list(sl, bullets,
                x + Inches(0.1), Inches(2.0), col_w - Inches(0.2), Inches(4.8),
                font_size=13)

# ── 8. Phase 2 — S5–S7 ────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Phase 2 — Server  ·  S5–S7")

# S5
card(sl, Inches(0.4), Inches(1.2), Inches(12.4), Inches(1.55))
txb(sl, "S5  ·  Postman Collection",
    Inches(0.6), Inches(1.28), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
txb(sl, "3 folders (Sessions, Players, QR Code)  ·  10 requests  ·  Test scripts auto-capture sessionId + playerId\n"
        "Status-code assertions on every request  ·  Enables full manual integration test without a browser",
    Inches(0.6), Inches(1.72), Inches(12.0), Inches(0.85), font_size=14, color=WHITE)

# S6
card(sl, Inches(0.4), Inches(2.9), Inches(12.4), Inches(1.9))
txb(sl, "S6  ·  HeartbeatMonitor",
    Inches(0.6), Inches(2.98), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
txb(sl, "BackgroundService running every 5 s\n"
        "Marks players stale after 30 s without heartbeat → broadcasts PlayerWentStale\n"
        "Ends sessions idle for 2+ hours → broadcasts GameEnded\n"
        "Recovery (PlayerReturned) handled in PlayersController on next push",
    Inches(0.6), Inches(3.42), Inches(12.0), Inches(1.2), font_size=14, color=WHITE)

# S7
card(sl, Inches(0.4), Inches(4.95), Inches(12.4), Inches(2.1))
txb(sl, "S7  ·  Web UI",
    Inches(0.6), Inches(5.03), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
bullet_list(sl, [
    "lobby.html + lobby.js — create session, QR display, player list, Start Spel button",
    "scoreboard.html + scoreboard.js — live player tiles, Beëindig spel, Nieuw Spel flow",
    "style.css — dark red #B71C1C background, orange #FFA726 accents, matching app palette",
    "App download QR on landing page (Google Play internal test link)",
    "signalr.min.js — local copy, no CDN dependency",
], Inches(0.6), Inches(5.47), Inches(12.0), Inches(1.5), font_size=14)

# ── 9. Phase 2 — D1–D3 Delphi Client ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Phase 2 — Delphi Client  ·  D1–D3")

# D1
card(sl, Inches(0.4), Inches(1.2), Inches(12.4), Inches(2.0))
txb(sl, "D1  ·  ServerClient.pas",
    Inches(0.6), Inches(1.28), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
txb(sl, "IServerClient interface: JoinSession, LeaveSession, PushScore, IsConnected\n"
        "TServerClient: THTTPClient per request · heartbeat TTimer (15 s) · offline after 3 consecutive failures\n"
        "TServerAwareScoreManager decorator: wraps IScoreManager, pushes score after every mutation\n"
        "Zero changes to existing MainFrm event handlers",
    Inches(0.6), Inches(1.72), Inches(12.0), Inches(1.35), font_size=14, color=WHITE)

# D2
card(sl, Inches(0.4), Inches(3.35), Inches(12.4), Inches(1.5))
txb(sl, "D2  ·  Menu Layout (MainFrm.fmx)",
    Inches(0.6), Inches(3.43), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
txb(sl, "pnlMenu height 150 → 250  ·  Added mnuJoinGame \"Aanmelden bij spel\"\n"
        "Added mnuLeaveGame \"Afmelden\" (initially hidden)  ·  lblStatus top-right: \"Online\" / \"Offline\" / empty",
    Inches(0.6), Inches(3.87), Inches(12.0), Inches(0.85), font_size=14, color=WHITE)

# D3
card(sl, Inches(0.4), Inches(5.0), Inches(12.4), Inches(2.0))
txb(sl, "D3  ·  Join / Leave Logic (MainFrm.pas)",
    Inches(0.6), Inches(5.08), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
txb(sl, "FServerClient: IServerClient initialised in FormCreate\n"
        "DoJoin: TThread.CreateAnonymousThread → JoinSession → swaps FScoreManager to TServerAwareScoreManager\n"
        "mnuLeaveGameClick: LeaveSession → swaps back to plain TScoreManager\n"
        "tmrStatusPoll (5 s): updates lblStatus from IsConnected on main thread",
    Inches(0.6), Inches(5.52), Inches(12.0), Inches(1.35), font_size=14, color=WHITE)

# ── 10. Phase 2 — D4–D5 ───────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Phase 2 — Delphi Client  ·  D4–D5")

card(sl, Inches(0.4), Inches(1.2), Inches(12.4), Inches(2.7))
txb(sl, "D4  ·  ZXing QR Scanner",
    Inches(0.6), Inches(1.28), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
bullet_list(sl, [
    "TScannerForm — full-screen camera preview, scans every 15th frame via TScanManager",
    "TOnResultProc named type (not TProc<string>) — avoids compiler overload resolution issues",
    "{$R *.fmx} requires ScannerFrm.fmx stub — omitting it raises EResNotFound at runtime",
    "Camera NOT deactivated in OnClose — deactivating inside DoFormClose hangs ZXing",
    "res/network_security_config.xml — Android 9+ blocks HTTP without it",
    "INTERNET + ACCESS_NETWORK_STATE via Delphi project options (not hardcoded in manifest)",
], Inches(0.6), Inches(1.72), Inches(12.0), Inches(2.0), font_size=13)

card(sl, Inches(0.4), Inches(4.05), Inches(12.4), Inches(2.9))
txb(sl, "D5  ·  End-to-End Test",
    Inches(0.6), Inches(4.13), Inches(12.0), Inches(0.45),
    font_size=17, bold=True, color=ORANGE)
bullet_list(sl, [
    "Full flow confirmed: browser creates session → player scans QR → joins → \"Online\" shown → scores update live",
    "Offline test: Wi-Fi off → app continues working → lblStatus shows \"Offline\" within ~45 s",
    "Recovery confirmed: Wi-Fi back → lblStatus returns to \"Online\" on next successful push",
    "VPN discovery: server bound to VPN interface (10.x.x.x) instead of LAN (192.168.x.x)",
    "Fix: IpAddressHelper now prefers 192.168.x.x > 172.16.x.x > 10.x.x.x",
], Inches(0.6), Inches(4.57), Inches(12.0), Inches(2.2), font_size=13)

# ── 11. Architecture ──────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Architecture Overview")

# Browser box
add_rect(sl, Inches(0.5), Inches(1.3), Inches(3.9), Inches(1.6), CARD)
txb(sl, "Browser", Inches(0.6), Inches(1.35), Inches(3.7), Inches(0.5),
    font_size=16, bold=True, color=ORANGE)
txb(sl, "lobby.html\nscoreboard.html",
    Inches(0.6), Inches(1.82), Inches(3.7), Inches(0.9), font_size=14, color=WHITE)

# Arrow: Browser ← SignalR ← Server
txb(sl, "◄── SignalR ──────────────────────",
    Inches(4.5), Inches(1.9), Inches(4.2), Inches(0.5),
    font_size=13, color=ORANGE, italic=True)

# Server box (centre-right)
add_rect(sl, Inches(8.8), Inches(1.3), Inches(4.1), Inches(4.8), CARD)
txb(sl, "ASP.NET Core 8 Server",
    Inches(8.95), Inches(1.35), Inches(3.9), Inches(0.5),
    font_size=16, bold=True, color=ORANGE)
bullet_list(sl, [
    "SessionStore",
    "(ConcurrentDictionary)",
    "",
    "HeartbeatMonitor",
    "(BackgroundService)",
    "",
    "GameHub (SignalR)",
    "",
    "REST Controllers",
], Inches(8.95), Inches(1.85), Inches(3.9), Inches(4.0), font_size=13)

# Android box
add_rect(sl, Inches(0.5), Inches(3.5), Inches(3.9), Inches(2.8), CARD)
txb(sl, "Delphi Android App",
    Inches(0.6), Inches(3.55), Inches(3.7), Inches(0.5),
    font_size=16, bold=True, color=ORANGE)
bullet_list(sl, [
    "IScoreManager",
    "  └─ TServerAwareScoreManager",
    "      └─ TScoreManager",
    "",
    "IServerClient",
    "  └─ TServerClient",
    "      TTimer (heartbeat 15 s)",
], Inches(0.6), Inches(4.05), Inches(3.7), Inches(2.1), font_size=13)

# Arrow: App → REST → Server
txb(sl, "────── REST ──────────────►",
    Inches(4.5), Inches(4.5), Inches(4.2), Inches(0.5),
    font_size=13, color=ORANGE, italic=True)

# Bottom note
txb(sl, "App → REST only.  Browser → SignalR only.  Server bridges both.",
    Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.45),
    font_size=14, color=DIM, italic=True, align=PP_ALIGN.CENTER)

# ── 12. Key Design Decisions ──────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Key Design Decisions")

decisions = [
    ("Interface + Decorator Pattern",
     "IScoreManager introduced in Phase 1 made the multiplayer extension zero-impact. "
     "TServerAwareScoreManager wraps the inner manager — no changes to any existing handler."),
    ("Offline-First",
     "TServerClient tracks a failure counter. After 3 consecutive failures IsConnected → false "
     "and SyncScore skips the push. The app degrades gracefully without user intervention."),
    ("In-Memory Session Store",
     "No database. ConcurrentDictionary is sufficient for a LAN game with a handful of players "
     "and a 2-hour session lifetime. Simplicity over persistence."),
    ("SignalR for Browser, REST for App",
     "Browser needs real-time push → SignalR. App only sends, never receives → plain HTTP POST "
     "on score change is simpler and avoids a WebSocket dependency in Delphi."),
    ("Unambiguous Session Codes",
     "Charset ABCDEFGHJKLMNPQRSTUVWXYZ23456789 — no 0/O or 1/I. "
     "6 characters, ~1.07 billion combinations. No collision risk at LAN scale."),
]

for i, (title, desc) in enumerate(decisions):
    y = Inches(1.2) + i * Inches(1.18)
    add_rect(sl, Inches(0.4), y + Inches(0.12), Inches(0.08), Inches(0.65), ORANGE)
    txb(sl, title,
        Inches(0.65), y, Inches(12.1), Inches(0.45),
        font_size=16, bold=True, color=ORANGE)
    txb(sl, desc,
        Inches(0.65), y + Inches(0.44), Inches(12.1), Inches(0.65),
        font_size=14, color=WHITE)

# ── 13. Notable Bugs ──────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "Notable Bugs Found and Fixed")

headers = ["Bug", "Symptom", "Fix"]
rows = [
    ["ParseJoinUrl outside try/except",
     "Silent failure on bad URL — TTask swallows exception",
     "Moved inside try block"],
    ["TTask.Run + Synchronize(nil,...)",
     "UI callback never fires on Android",
     "TThread.CreateAnonymousThread + Synchronize(CurrentThread,...)"],
    ["Camera deactivation in OnClose",
     "ZXing hangs, app freezes",
     "Removed Active := False from DoFormClose"],
    ["Missing ScannerFrm.fmx stub",
     "EResNotFound at runtime",
     "Created minimal FMX stub file"],
    ["VPN interface selected as LAN IP",
     "QR encodes unreachable 10.x.x.x address",
     "IpAddressHelper prefers 192.168.x.x"],
    ["End Spel button did nothing",
     "Relied on SignalR event that could be dropped",
     "Call showGameEnded() directly on successful DELETE"],
]
col_widths = [Inches(3.4), Inches(4.5), Inches(4.5)]
table_slide(sl, headers, rows,
            left=Inches(0.4), top=Inches(1.2), width=Inches(12.4),
            col_widths=col_widths, header_size=15, row_size=13)

# ── 14. REST API Reference ────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)
heading(sl, "REST API Reference")

headers = ["Method", "Endpoint", "Purpose"]
rows = [
    ["POST",   "/api/sessions",                               "Create session → { sessionId, joinUrl }"],
    ["GET",    "/api/sessions/{id}",                          "Get session state + player list"],
    ["POST",   "/api/sessions/{id}/start",                    "Start the game"],
    ["DELETE", "/api/sessions/{id}",                          "End and delete session → GameEnded"],
    ["POST",   "/api/sessions/{id}/players",                  "Register player → { playerId }"],
    ["GET",    "/api/sessions/{id}/players",                  "List all players"],
    ["PUT",    "/api/sessions/{id}/players/{pid}/score",      "Push score update → ScoreUpdated"],
    ["POST",   "/api/sessions/{id}/players/{pid}/heartbeat",  "Keep-alive (every 15 s)"],
    ["GET",    "/api/sessions/{id}/qr",                       "QR code PNG encoding the join URL"],
    ["GET",    "/api/appqr",                                  "QR code PNG for app download link"],
]
col_widths = [Inches(1.1), Inches(5.5), Inches(5.8)]
table_slide(sl, headers, rows,
            left=Inches(0.4), top=Inches(1.15), width=Inches(12.4),
            col_widths=col_widths, header_size=14, row_size=13)

# ── 15. Final slide ───────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_background(sl)

for x, y, s in [
    (Inches(10.5), Inches(-0.5), Inches(4)),
    (Inches(11.5), Inches(4.5),  Inches(3.5)),
    (Inches(-0.5), Inches(5.0),  Inches(2.5)),
]:
    add_circle(sl, x, y, s, CARD)

add_rect(sl, 0, Inches(2.8), Inches(0.18), Inches(2.2), ORANGE)

txb(sl, "Project Complete",
    Inches(0.5), Inches(2.8), Inches(11), Inches(1.1),
    font_size=52, bold=True, color=WHITE)

txb(sl, "De Slimste Mens Timer",
    Inches(0.5), Inches(3.95), Inches(11), Inches(0.65),
    font_size=26, bold=True, color=ORANGE)

bullet_list(sl, [
    "Delphi / FireMonkey Android app — 9 steps, Phase 1",
    "ASP.NET Core 8 server + real-time SignalR scoreboard — Phase 2",
    "Each step proposed, reviewed, and approved before execution",
    "Built entirely with Claude Code",
], Inches(0.5), Inches(4.65), Inches(11.5), Inches(2.0), font_size=18)

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
